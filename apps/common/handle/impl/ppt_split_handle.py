# import os
# from typing import List
#
# # from langchain_community.document_loaders.unstructured import UnstructuredFileLoader
#
#
#
# """Loader that uses unstructured to load files."""
#
# import collections
# from abc import ABC, abstractmethod
# from pathlib import Path
# from typing import IO, Any, Callable, Dict, Iterator, List, Optional, Sequence, Union
#
# from langchain_core.documents import Document
# """Abstract interface for document loader implementations."""
#
# from __future__ import annotations
#
# from abc import ABC, abstractmethod
# from typing import TYPE_CHECKING, AsyncIterator, Iterator, List, Optional
#
# from langchain_core.documents import Document
# from langchain_core.runnables import run_in_executor
#
# if TYPE_CHECKING:
#     from langchain_text_splitters import TextSplitter
#
# import contextlib
# import mimetypes
# from io import BufferedReader, BytesIO
# from pathlib import PurePath
# from typing import Any, Generator, List, Literal, Mapping, Optional, Union, cast
#
# from langchain_core.load.serializable import Serializable
# from langchain_core.pydantic_v1 import Field, root_validator
#
# PathLike = Union[str, PurePath]
#
#
# class BaseMedia(Serializable):
#     """Use to represent media content.
#
#     Media objets can be used to represent raw data, such as text or binary data.
#
#     LangChain Media objects allow associating metadata and an optional identifier
#     with the content.
#
#     The presence of an ID and metadata make it easier to store, index, and search
#     over the content in a structured way.
#     """
#
#     # The ID field is optional at the moment.
#     # It will likely become required in a future major release after
#     # it has been adopted by enough vectorstore implementations.
#     id: Optional[str] = None
#     """An optional identifier for the document.
#
#     Ideally this should be unique across the document collection and formatted
#     as a UUID, but this will not be enforced.
#
#     .. versionadded:: 0.2.11
#     """
#
#     metadata: dict = Field(default_factory=dict)
#     """Arbitrary metadata associated with the content."""
#
#
# class Blob(BaseMedia):
#     """Blob represents raw data by either reference or value.
#
#     Provides an interface to materialize the blob in different representations, and
#     help to decouple the development of data loaders from the downstream parsing of
#     the raw data.
#
#     Inspired by: https://developer.mozilla.org/en-US/docs/Web/API/Blob
#
#     Example: Initialize a blob from in-memory data
#
#         .. code-block:: python
#
#             from langchain_core.documents import Blob
#
#             blob = Blob.from_data("Hello, world!")
#
#             # Read the blob as a string
#             print(blob.as_string())
#
#             # Read the blob as bytes
#             print(blob.as_bytes())
#
#             # Read the blob as a byte stream
#             with blob.as_bytes_io() as f:
#                 print(f.read())
#
#     Example: Load from memory and specify mime-type and metadata
#
#         .. code-block:: python
#
#             from langchain_core.documents import Blob
#
#             blob = Blob.from_data(
#                 data="Hello, world!",
#                 mime_type="text/plain",
#                 metadata={"source": "https://example.com"}
#             )
#
#     Example: Load the blob from a file
#
#         .. code-block:: python
#
#             from langchain_core.documents import Blob
#
#             blob = Blob.from_path("path/to/file.txt")
#
#             # Read the blob as a string
#             print(blob.as_string())
#
#             # Read the blob as bytes
#             print(blob.as_bytes())
#
#             # Read the blob as a byte stream
#             with blob.as_bytes_io() as f:
#                 print(f.read())
#     """
#
#     data: Union[bytes, str, None]
#     """Raw data associated with the blob."""
#     mimetype: Optional[str] = None
#     """MimeType not to be confused with a file extension."""
#     encoding: str = "utf-8"
#     """Encoding to use if decoding the bytes into a string.
#
#     Use utf-8 as default encoding, if decoding to string.
#     """
#     path: Optional[PathLike] = None
#     """Location where the original content was found."""
#
#     class Config:
#         arbitrary_types_allowed = True
#         frozen = True
#
#     @property
#     def source(self) -> Optional[str]:
#         """The source location of the blob as string if known otherwise none.
#
#         If a path is associated with the blob, it will default to the path location.
#
#         Unless explicitly set via a metadata field called "source", in which
#         case that value will be used instead.
#         """
#         if self.metadata and "source" in self.metadata:
#             return cast(Optional[str], self.metadata["source"])
#         return str(self.path) if self.path else None
#
#     @root_validator(pre=True)
#     def check_blob_is_valid(cls, values: Mapping[str, Any]) -> Mapping[str, Any]:
#         """Verify that either data or path is provided."""
#         if "data" not in values and "path" not in values:
#             raise ValueError("Either data or path must be provided")
#         return values
#
#     def as_string(self) -> str:
#         """Read data as a string."""
#         if self.data is None and self.path:
#             with open(str(self.path), "r", encoding=self.encoding) as f:
#                 return f.read()
#         elif isinstance(self.data, bytes):
#             return self.data.decode(self.encoding)
#         elif isinstance(self.data, str):
#             return self.data
#         else:
#             raise ValueError(f"Unable to get string for blob {self}")
#
#     def as_bytes(self) -> bytes:
#         """Read data as bytes."""
#         if isinstance(self.data, bytes):
#             return self.data
#         elif isinstance(self.data, str):
#             return self.data.encode(self.encoding)
#         elif self.data is None and self.path:
#             with open(str(self.path), "rb") as f:
#                 return f.read()
#         else:
#             raise ValueError(f"Unable to get bytes for blob {self}")
#
#     @contextlib.contextmanager
#     def as_bytes_io(self) -> Generator[Union[BytesIO, BufferedReader], None, None]:
#         """Read data as a byte stream."""
#         if isinstance(self.data, bytes):
#             yield BytesIO(self.data)
#         elif self.data is None and self.path:
#             with open(str(self.path), "rb") as f:
#                 yield f
#         else:
#             raise NotImplementedError(f"Unable to convert blob {self}")
#
#     @classmethod
#     def from_path(
#             cls,
#             path: PathLike,
#             *,
#             encoding: str = "utf-8",
#             mime_type: Optional[str] = None,
#             guess_type: bool = True,
#             metadata: Optional[dict] = None,
#     ) -> Blob:
#         """Load the blob from a path like object.
#
#         Args:
#             path: path like object to file to be read
#             encoding: Encoding to use if decoding the bytes into a string
#             mime_type: if provided, will be set as the mime-type of the data
#             guess_type: If True, the mimetype will be guessed from the file extension,
#                         if a mime-type was not provided
#             metadata: Metadata to associate with the blob
#
#         Returns:
#             Blob instance
#         """
#         if mime_type is None and guess_type:
#             _mimetype = mimetypes.guess_type(path)[0] if guess_type else None
#         else:
#             _mimetype = mime_type
#         # We do not load the data immediately, instead we treat the blob as a
#         # reference to the underlying data.
#         return cls(
#             data=None,
#             mimetype=_mimetype,
#             encoding=encoding,
#             path=path,
#             metadata=metadata if metadata is not None else {},
#         )
#
#     @classmethod
#     def from_data(
#             cls,
#             data: Union[str, bytes],
#             *,
#             encoding: str = "utf-8",
#             mime_type: Optional[str] = None,
#             path: Optional[str] = None,
#             metadata: Optional[dict] = None,
#     ) -> Blob:
#         """Initialize the blob from in-memory data.
#
#         Args:
#             data: the in-memory data associated with the blob
#             encoding: Encoding to use if decoding the bytes into a string
#             mime_type: if provided, will be set as the mime-type of the data
#             path: if provided, will be set as the source from which the data came
#             metadata: Metadata to associate with the blob
#
#         Returns:
#             Blob instance
#         """
#         return cls(
#             data=data,
#             mimetype=mime_type,
#             encoding=encoding,
#             path=path,
#             metadata=metadata if metadata is not None else {},
#         )
#
#     def __repr__(self) -> str:
#         """Define the blob representation."""
#         str_repr = f"Blob {id(self)}"
#         if self.source:
#             str_repr += f" {self.source}"
#         return str_repr
#
#
# class Document(BaseMedia):
#     """Class for storing a piece of text and associated metadata.
#
#     Example:
#
#         .. code-block:: python
#
#             from langchain_core.documents import Document
#
#             document = Document(
#                 page_content="Hello, world!",
#                 metadata={"source": "https://example.com"}
#             )
#     """
#
#     page_content: str
#     """String text."""
#     type: Literal["Document"] = "Document"
#
#     def __init__(self, page_content: str, **kwargs: Any) -> None:
#         """Pass page_content in as positional or named arg."""
#         # my-py is complaining that page_content is not defined on the base class.
#         # Here, we're relying on pydantic base class to handle the validation.
#         super().__init__(page_content=page_content, **kwargs)  # type: ignore[call-arg]
#
#     @classmethod
#     def is_lc_serializable(cls) -> bool:
#         """Return whether this class is serializable."""
#         return True
#
#     @classmethod
#     def get_lc_namespace(cls) -> List[str]:
#         """Get the namespace of the langchain object."""
#         return ["langchain", "schema", "document"]
#
#     def __str__(self) -> str:
#         """Override __str__ to restrict it to page_content and metadata."""
#         # The format matches pydantic format for __str__.
#         #
#         # The purpose of this change is to make sure that user code that
#         # feeds Document objects directly into prompts remains unchanged
#         # due to the addition of the id field (or any other fields in the future).
#         #
#         # This override will likely be removed in the future in favor of
#         # a more general solution of formatting content directly inside the prompts.
#         if self.metadata:
#             return f"page_content='{self.page_content}' metadata={self.metadata}"
#         else:
#             return f"page_content='{self.page_content}'"
#
#
# class BaseLoader(ABC):
#     """Interface for Document Loader.
#
#     Implementations should implement the lazy-loading method using generators
#     to avoid loading all Documents into memory at once.
#
#     `load` is provided just for user convenience and should not be overridden.
#     """
#
#     # Sub-classes should not implement this method directly. Instead, they
#     # should implement the lazy load method.
#     def load(self) -> List[Document]:
#         """Load data into Document objects."""
#         return list(self.lazy_load())
#
#     async def aload(self) -> List[Document]:
#         """Load data into Document objects."""
#         return [document async for document in self.alazy_load()]
#
#     def load_and_split(
#         self, text_splitter: Optional[TextSplitter] = None
#     ) -> List[Document]:
#         """Load Documents and split into chunks. Chunks are returned as Documents.
#
#         Do not override this method. It should be considered to be deprecated!
#
#         Args:
#             text_splitter: TextSplitter instance to use for splitting documents.
#               Defaults to RecursiveCharacterTextSplitter.
#
#         Returns:
#             List of Documents.
#         """
#
#         if text_splitter is None:
#             try:
#                 from langchain_text_splitters import RecursiveCharacterTextSplitter
#             except ImportError as e:
#                 raise ImportError(
#                     "Unable to import from langchain_text_splitters. Please specify "
#                     "text_splitter or install langchain_text_splitters with "
#                     "`pip install -U langchain-text-splitters`."
#                 ) from e
#
#             _text_splitter: TextSplitter = RecursiveCharacterTextSplitter()
#         else:
#             _text_splitter = text_splitter
#         docs = self.load()
#         return _text_splitter.split_documents(docs)
#
#     # Attention: This method will be upgraded into an abstractmethod once it's
#     #            implemented in all the existing subclasses.
#     def lazy_load(self) -> Iterator[Document]:
#         """A lazy loader for Documents."""
#         if type(self).load != BaseLoader.load:
#             return iter(self.load())
#         raise NotImplementedError(
#             f"{self.__class__.__name__} does not implement lazy_load()"
#         )
#
#     async def alazy_load(self) -> AsyncIterator[Document]:
#         """A lazy loader for Documents."""
#         iterator = await run_in_executor(None, self.lazy_load)
#         done = object()
#         while True:
#             doc = await run_in_executor(None, next, iterator, done)  # type: ignore[call-arg, arg-type]
#             if doc is done:
#                 break
#             yield doc  # type: ignore[misc]
#
#
# class BaseBlobParser(ABC):
#     """Abstract interface for blob parsers.
#
#     A blob parser provides a way to parse raw data stored in a blob into one
#     or more documents.
#
#     The parser can be composed with blob loaders, making it easy to reuse
#     a parser independent of how the blob was originally loaded.
#     """
#
#     @abstractmethod
#     def lazy_parse(self, blob: Blob) -> Iterator[Document]:
#         """Lazy parsing interface.
#
#         Subclasses are required to implement this method.
#
#         Args:
#             blob: Blob instance
#
#         Returns:
#             Generator of documents
#         """
#
#     def parse(self, blob: Blob) -> List[Document]:
#         """Eagerly parse the blob into a document or documents.
#
#         This is a convenience method for interactive development environment.
#
#         Production applications should favor the lazy_parse method instead.
#
#         Subclasses should generally not over-ride this parse method.
#
#         Args:
#             blob: Blob instance
#
#         Returns:
#             List of documents
#         """
#         return list(self.lazy_parse(blob))
#
#
# def satisfies_min_unstructured_version(min_version: str) -> bool:
#     """Check if the installed `Unstructured` version exceeds the minimum version
#     for the feature in question."""
#     from unstructured.__version__ import __version__ as __unstructured_version__
#
#     min_version_tuple = tuple([int(x) for x in min_version.split(".")])
#
#     # NOTE(MthwRobinson) - enables the loader to work when you're using pre-release
#     # versions of unstructured like 0.4.17-dev1
#     _unstructured_version = __unstructured_version__.split("-")[0]
#     unstructured_version_tuple = tuple(
#         [int(x) for x in _unstructured_version.split(".")]
#     )
#
#     return unstructured_version_tuple >= min_version_tuple
#
#
# def validate_unstructured_version(min_unstructured_version: str) -> None:
#     """Raise an error if the `Unstructured` version does not exceed the
#     specified minimum."""
#     if not satisfies_min_unstructured_version(min_unstructured_version):
#         raise ValueError(
#             f"unstructured>={min_unstructured_version} is required in this loader."
#         )
#
#
# class UnstructuredBaseLoader(BaseLoader, ABC):
#     """Base Loader that uses `Unstructured`."""
#
#     def __init__(
#         self,
#         mode: str = "single",
#         post_processors: Optional[List[Callable]] = None,
#         **unstructured_kwargs: Any,
#     ):
#         """Initialize with file path."""
#         try:
#             import unstructured  # noqa:F401
#         except ImportError:
#             raise ImportError(
#                 "unstructured package not found, please install it with "
#                 "`pip install unstructured`"
#             )
#         _valid_modes = {"single", "elements", "paged"}
#         if mode not in _valid_modes:
#             raise ValueError(
#                 f"Got {mode} for `mode`, but should be one of `{_valid_modes}`"
#             )
#         self.mode = mode
#
#         if not satisfies_min_unstructured_version("0.5.4"):
#             if "strategy" in unstructured_kwargs:
#                 unstructured_kwargs.pop("strategy")
#
#         self.unstructured_kwargs = unstructured_kwargs
#         self.post_processors = post_processors or []
#
#     @abstractmethod
#     def _get_elements(self) -> List:
#         """Get elements."""
#
#     @abstractmethod
#     def _get_metadata(self) -> dict:
#         """Get metadata."""
#
#     def _post_process_elements(self, elements: list) -> list:
#         """Applies post processing functions to extracted unstructured elements.
#         Post processing functions are str -> str callables are passed
#         in using the post_processors kwarg when the loader is instantiated."""
#         for element in elements:
#             for post_processor in self.post_processors:
#                 element.apply(post_processor)
#         return elements
#
#     def lazy_load(self) -> Iterator[Document]:
#         """Load file."""
#         elements = self._get_elements()
#         self._post_process_elements(elements)
#         if self.mode == "elements":
#             for element in elements:
#                 metadata = self._get_metadata()
#                 # NOTE(MthwRobinson) - the attribute check is for backward compatibility
#                 # with unstructured<0.4.9. The metadata attributed was added in 0.4.9.
#                 if hasattr(element, "metadata"):
#                     metadata.update(element.metadata.to_dict())
#                 if hasattr(element, "category"):
#                     metadata["category"] = element.category
#                 yield Document(page_content=str(element), metadata=metadata)
#         elif self.mode == "paged":
#             text_dict: Dict[int, str] = {}
#             meta_dict: Dict[int, Dict] = {}
#
#             for idx, element in enumerate(elements):
#                 metadata = self._get_metadata()
#                 if hasattr(element, "metadata"):
#                     metadata.update(element.metadata.to_dict())
#                 page_number = metadata.get("page_number", 1)
#
#                 # Check if this page_number already exists in docs_dict
#                 if page_number not in text_dict:
#                     # If not, create new entry with initial text and metadata
#                     text_dict[page_number] = str(element) + "\n\n"
#                     meta_dict[page_number] = metadata
#                 else:
#                     # If exists, append to text and update the metadata
#                     text_dict[page_number] += str(element) + "\n\n"
#                     meta_dict[page_number].update(metadata)
#
#             # Convert the dict to a list of Document objects
#             for key in text_dict.keys():
#                 yield Document(page_content=text_dict[key], metadata=meta_dict[key])
#         elif self.mode == "single":
#             metadata = self._get_metadata()
#             text = "\n\n".join([str(el) for el in elements])
#             yield Document(page_content=text, metadata=metadata)
#         else:
#             raise ValueError(f"mode of {self.mode} not supported.")
#
#
# class UnstructuredFileLoader(UnstructuredBaseLoader):
#     """Load files using `Unstructured`.
#
#     The file loader uses the
#     unstructured partition function and will automatically detect the file
#     type. You can run the loader in one of two modes: "single" and "elements".
#     If you use "single" mode, the document will be returned as a single
#     langchain Document object. If you use "elements" mode, the unstructured
#     library will split the document into elements such as Title and NarrativeText.
#     You can pass in additional unstructured kwargs after mode to apply
#     different unstructured settings.
#
#     Examples
#     --------
#     from langchain_community.document_loaders import UnstructuredFileLoader
#
#     loader = UnstructuredFileLoader(
#         "example.pdf", mode="elements", strategy="fast",
#     )
#     docs = loader.load()
#
#     References
#     ----------
#     https://unstructured-io.github.io/unstructured/bricks.html#partition
#     """
#
#     def __init__(
#         self,
#         file_path: Union[str, List[str], Path, List[Path], None],
#         mode: str = "single",
#         **unstructured_kwargs: Any,
#     ):
#         """Initialize with file path."""
#         self.file_path = file_path
#         super().__init__(mode=mode, **unstructured_kwargs)
#
#     def _get_elements(self) -> List:
#         from unstructured.partition.auto import partition
#
#         if isinstance(self.file_path, list):
#             elements = []
#             for file in self.file_path:
#                 if isinstance(file, Path):
#                     file = str(file)
#                 elements.extend(partition(filename=file, **self.unstructured_kwargs))
#             return elements
#         else:
#             if isinstance(self.file_path, Path):
#                 self.file_path = str(self.file_path)
#             return partition(filename=self.file_path, **self.unstructured_kwargs)
#
#     def _get_metadata(self) -> dict:
#         return {"source": self.file_path}
#
#
# def get_elements_from_api(
#     file_path: Union[str, List[str], Path, List[Path], None] = None,
#     file: Union[IO, Sequence[IO], None] = None,
#     api_url: str = "https://api.unstructured.io/general/v0/general",
#     api_key: str = "",
#     **unstructured_kwargs: Any,
# ) -> List:
#     """Retrieve a list of elements from the `Unstructured API`."""
#     if is_list := isinstance(file_path, list):
#         file_path = [str(path) for path in file_path]
#     if isinstance(file, collections.abc.Sequence) or is_list:
#         from unstructured.partition.api import partition_multiple_via_api
#
#         _doc_elements = partition_multiple_via_api(
#             filenames=file_path,
#             files=file,
#             api_key=api_key,
#             api_url=api_url,
#             **unstructured_kwargs,
#         )
#
#         elements = []
#         for _elements in _doc_elements:
#             elements.extend(_elements)
#
#         return elements
#     else:
#         from unstructured.partition.api import partition_via_api
#
#         return partition_via_api(
#             filename=str(file_path) if file_path is not None else None,
#             file=file,
#             api_key=api_key,
#             api_url=api_url,
#             **unstructured_kwargs,
#         )
#
#
# class UnstructuredAPIFileLoader(UnstructuredFileLoader):
#     """Load files using `Unstructured` API.
#
#     By default, the loader makes a call to the hosted Unstructured API.
#     If you are running the unstructured API locally, you can change the
#     API rule by passing in the url parameter when you initialize the loader.
#     The hosted Unstructured API requires an API key. See
#     https://www.unstructured.io/api-key/ if you need to generate a key.
#
#     You can run the loader in one of two modes: "single" and "elements".
#     If you use "single" mode, the document will be returned as a single
#     langchain Document object. If you use "elements" mode, the unstructured
#     library will split the document into elements such as Title and NarrativeText.
#     You can pass in additional unstructured kwargs after mode to apply
#     different unstructured settings.
#
#     Examples
#     ```python
#     from langchain_community.document_loaders import UnstructuredAPIFileLoader
#
#     loader = UnstructuredFileAPILoader(
#         "example.pdf", mode="elements", strategy="fast", api_key="MY_API_KEY",
#     )
#     docs = loader.load()
#
#     References
#     ----------
#     https://unstructured-io.github.io/unstructured/bricks.html#partition
#     https://www.unstructured.io/api-key/
#     https://github.com/Unstructured-IO/unstructured-api
#     """
#
#     def __init__(
#         self,
#         file_path: Union[str, List[str], None] = "",
#         mode: str = "single",
#         url: str = "https://api.unstructured.io/general/v0/general",
#         api_key: str = "",
#         **unstructured_kwargs: Any,
#     ):
#         """Initialize with file path."""
#
#         validate_unstructured_version(min_unstructured_version="0.10.15")
#
#         self.url = url
#         self.api_key = api_key
#
#         super().__init__(file_path=file_path, mode=mode, **unstructured_kwargs)
#
#     def _get_metadata(self) -> dict:
#         return {"source": self.file_path}
#
#     def _get_elements(self) -> List:
#         return get_elements_from_api(
#             file_path=self.file_path,
#             api_key=self.api_key,
#             api_url=self.url,
#             **self.unstructured_kwargs,
#         )
#
#
# class UnstructuredFileIOLoader(UnstructuredBaseLoader):
#     """Load files using `Unstructured`.
#
#     The file loader
#     uses the unstructured partition function and will automatically detect the file
#     type. You can run the loader in one of two modes: "single" and "elements".
#     If you use "single" mode, the document will be returned as a single
#     langchain Document object. If you use "elements" mode, the unstructured
#     library will split the document into elements such as Title and NarrativeText.
#     You can pass in additional unstructured kwargs after mode to apply
#     different unstructured settings.
#
#     Examples
#     --------
#     from langchain_community.document_loaders import UnstructuredFileIOLoader
#
#     with open("example.pdf", "rb") as f:
#         loader = UnstructuredFileIOLoader(
#             f, mode="elements", strategy="fast",
#         )
#         docs = loader.load()
#
#
#     References
#     ----------
#     https://unstructured-io.github.io/unstructured/bricks.html#partition
#     """
#
#     def __init__(
#         self,
#         file: Union[IO, Sequence[IO]],
#         mode: str = "single",
#         **unstructured_kwargs: Any,
#     ):
#         """Initialize with file path."""
#         self.file = file
#         super().__init__(mode=mode, **unstructured_kwargs)
#
#     def _get_elements(self) -> List:
#         from unstructured.partition.auto import partition
#
#         return partition(file=self.file, **self.unstructured_kwargs)
#
#     def _get_metadata(self) -> dict:
#         return {}
#
#
# class UnstructuredAPIFileIOLoader(UnstructuredFileIOLoader):
#     """Load files using `Unstructured` API.
#
#     By default, the loader makes a call to the hosted Unstructured API.
#     If you are running the unstructured API locally, you can change the
#     API rule by passing in the url parameter when you initialize the loader.
#     The hosted Unstructured API requires an API key. See
#     https://www.unstructured.io/api-key/ if you need to generate a key.
#
#     You can run the loader in one of two modes: "single" and "elements".
#     If you use "single" mode, the document will be returned as a single
#     langchain Document object. If you use "elements" mode, the unstructured
#     library will split the document into elements such as Title and NarrativeText.
#     You can pass in additional unstructured kwargs after mode to apply
#     different unstructured settings.
#
#     Examples
#     --------
#     from langchain_community.document_loaders import UnstructuredAPIFileLoader
#
#     with open("example.pdf", "rb") as f:
#         loader = UnstructuredFileAPILoader(
#             f, mode="elements", strategy="fast", api_key="MY_API_KEY",
#         )
#         docs = loader.load()
#
#     References
#     ----------
#     https://unstructured-io.github.io/unstructured/bricks.html#partition
#     https://www.unstructured.io/api-key/
#     https://github.com/Unstructured-IO/unstructured-api
#     """
#
#     def __init__(
#         self,
#         file: Union[IO, Sequence[IO]],
#         mode: str = "single",
#         url: str = "https://api.unstructured.io/general/v0/general",
#         api_key: str = "",
#         **unstructured_kwargs: Any,
#     ):
#         """Initialize with file path."""
#
#         if isinstance(file, collections.abc.Sequence):
#             validate_unstructured_version(min_unstructured_version="0.6.3")
#         if file:
#             validate_unstructured_version(min_unstructured_version="0.6.2")
#
#         self.url = url
#         self.api_key = api_key
#
#         super().__init__(file=file, mode=mode, **unstructured_kwargs)
#
#     def _get_elements(self) -> List:
#         return get_elements_from_api(
#             file=self.file,
#             api_key=self.api_key,
#             api_url=self.url,
#             **self.unstructured_kwargs,
#         )
#
# class UnstructuredPowerPointLoader(UnstructuredFileLoader):
#     """Load `Microsoft PowerPoint` files using `Unstructured`.
#
#     Works with both .ppt and .pptx files.
#     You can run the loader in one of two modes: "single" and "elements".
#     If you use "single" mode, the document will be returned as a single
#     langchain Document object. If you use "elements" mode, the unstructured
#     library will split the document into elements such as Title and NarrativeText.
#     You can pass in additional unstructured kwargs after mode to apply
#     different unstructured settings.
#
#     Examples
#     --------
#     from langchain_community.document_loaders import UnstructuredPowerPointLoader
#
#     loader = UnstructuredPowerPointLoader(
#         "example.pptx", mode="elements", strategy="fast",
#     )
#     docs = loader.load()
#
#     References
#     ----------
#     https://unstructured-io.github.io/unstructured/bricks.html#partition-pptx
#     """
#
#     def _get_elements(self) -> List:
#         from unstructured.__version__ import __version__ as __unstructured_version__
#         from unstructured.file_utils.filetype import FileType, detect_filetype
#
#         unstructured_version = tuple(
#             [int(x) for x in __unstructured_version__.split(".")]
#         )
#         # NOTE(MthwRobinson) - magic will raise an import error if the libmagic
#         # system dependency isn't installed. If it's not installed, we'll just
#         # check the file extension
#         try:
#             import magic  # noqa: F401
#
#             is_ppt = detect_filetype(self.file_path) == FileType.PPT
#         except ImportError:
#             _, extension = os.path.splitext(str(self.file_path))
#             is_ppt = extension == ".ppt"
#
#         if is_ppt and unstructured_version < (0, 4, 11):
#             raise ValueError(
#                 f"You are on unstructured version {__unstructured_version__}. "
#                 "Partitioning .ppt files is only supported in unstructured>=0.4.11. "
#                 "Please upgrade the unstructured package and try again."
#             )
#
#         if is_ppt:
#             from unstructured.partition.ppt import partition_ppt
#
#             return partition_ppt(filename=self.file_path, **self.unstructured_kwargs)
#         else:
#             from unstructured.partition.pptx import partition_pptx
#
#             return partition_pptx(filename=self.file_path, **self.unstructured_kwargs)
"""


from typing import List

import tqdm
from langchain_community.document_loaders.unstructured import UnstructuredFileLoader


class RapidOCRPPTLoader(UnstructuredFileLoader):
    def _get_elements(self) -> List:
        def ppt2text(filepath):
            from io import BytesIO

            import numpy as np
            from PIL import Image
            from pptx import Presentation
            from rapidocr_onnxruntime import RapidOCR

            ocr = RapidOCR()
            prs = Presentation(filepath)
            resp = ""

            def extract_text(shape):
                nonlocal resp
                if shape.has_text_frame:
                    resp += shape.text.strip() + "\n"
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            for paragraph in cell.text_frame.paragraphs:
                                resp += paragraph.text.strip() + "\n"
                if shape.shape_type == 13:  # 13 表示图片
                    image = Image.open(BytesIO(shape.image.blob))
                    result, _ = ocr(np.array(image))
                    if result:
                        ocr_result = [line[1] for line in result]
                        resp += "\n".join(ocr_result)
                elif shape.shape_type == 6:  # 6 表示组合
                    for child_shape in shape.shapes:
                        extract_text(child_shape)

            b_unit = tqdm.tqdm(
                total=len(prs.slides), desc="RapidOCRPPTLoader slide index: 1"
            )
            # 遍历所有幻灯片
            for slide_number, slide in enumerate(prs.slides, start=1):
                b_unit.set_description(
                    "RapidOCRPPTLoader slide index: {}".format(slide_number)
                )
                b_unit.refresh()
                sorted_shapes = sorted(
                    slide.shapes, key=lambda x: (x.top, x.left)
                )  # 从上到下、从左到右遍历
                for shape in sorted_shapes:
                    extract_text(shape)
                b_unit.update(1)
            return resp

        text = ppt2text(self.file_path)
        from unstructured.partition.text import partition_text

        return partition_text(text=text, **self.unstructured_kwargs)


if __name__ == "__main__":
    loader = RapidOCRPPTLoader(file_path="../tests/samples/ocr_test.pptx")
    docs = loader.load()
    print(docs)

"""

# 解压并提取 PPTX 文件中图片
import re
import traceback
import uuid
from operator import attrgetter
from typing import List

# from PIL import Image
from pptx.enum.shapes import MSO_SHAPE_TYPE

from common.handle.base_split_handle import BaseSplitHandle
from common.util.split_model import SplitModel
from dataset.models import Image

"""
    import os
    import zipfile


    def extract_images_from_pptx(pptx_file_path, output_dir):
        with zipfile.ZipFile(pptx_file_path, 'r') as zip_ref:
            for file in zip_ref.namelist():

                if file.startswith('ppt/media/') and file.endswith(('.png', '.jpg', '.jpeg')):
                    print(file)
                    zip_ref.extract(file, output_dir)
"""

default_pattern_list = [re.compile('(?<=^)# .*|(?<=\\n)# .*'),
                        re.compile('(?<=\\n)(?<!#)## (?!#).*|(?<=^)(?<!#)## (?!#).*'),
                        re.compile("(?<=\\n)(?<!#)### (?!#).*|(?<=^)(?<!#)### (?!#).*"),
                        re.compile("(?<=\\n)(?<!#)#### (?!#).*|(?<=^)(?<!#)#### (?!#).*"),
                        re.compile("(?<=\\n)(?<!#)##### (?!#).*|(?<=^)(?<!#)##### (?!#).*"),
                        re.compile("(?<=\\n)(?<!#)###### (?!#).*|(?<=^)(?<!#)###### (?!#).*")]


def ungroup_shapes(shapes):
    res = []
    for shape in shapes:
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                res.extend(ungroup_shapes(shape.shapes))
            else:
                res.append(shape)
        except Exception as e:
            print(f'failed to load shape {shape}, skipped. error: {e}')
    return res


class PPTSplitHandle(BaseSplitHandle):
    @staticmethod
    def table_to_md(rows):
        # 创建 Markdown 格式的表格
        md_table = '| ' + ' | '.join(
            [cell for cell in rows[0]]) + ' |\n'
        md_table += '| ' + ' | '.join(['---' for i in range(len(rows[0]))]) + ' |\n'
        for row in rows[1:]:
            md_table += '| ' + ' | '.join(
                [cell for cell in row]) + ' |\n'
        return md_table

    def to_md(self, cells):
        elements = []
        page = -1
        for item in cells:
            if item["type"] == "table":
                table = self.table_to_md(item["content"])
                elements.append(table)
            elif item["type"] == "image":
                elements.append(f'![](/api/image/{item["meta"]["id"]})')
            else:
                if page < item["page"]:
                    content = "# " + item["content"]
                else:
                    content = item["content"]
                elements.append(content)
                page = item["page"]
        return "\n".join(elements)

    @staticmethod
    def parse_ppt(file) -> List[dict]:
        from unstructured.partition.ppt import partition_ppt
        elements = partition_ppt(file=file)
        return []

    @staticmethod
    def parse_pptx(file) -> List[dict]:
        from pptx import Presentation
        prs = Presentation(file)
        cells = []
        for idx, slide in enumerate(prs.slides):
            try:
                shapes = sorted(ungroup_shapes(slide.shapes), key=attrgetter('top', 'left'))
            except BaseException as e:
                print(e)
                continue
            for shape in shapes:
                item = {"page": idx, "type": "", "shape_type": shape.shape_type, "content": "", "left": shape.left,
                        "top": shape.top, "height": shape.height, "width": shape.width}
                if shape.shape_type == 13:
                    image = shape.image.blob  # Image.open(BytesIO(shape.image.blob))
                    item["type"] = "image"
                    item["meta"] = {"id": shape.image.sha1, "dpi": shape.image.dpi, "name": shape.name}
                    item["content"] = image
                elif shape.has_table:
                    table = []
                    for row in shape.table.rows:
                        table.append(
                            [paragraph.text.strip() for cell in row.cells for paragraph in cell.text_frame.paragraphs])
                    item["type"] = "table"
                    item["content"] = table
                elif shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if not text:
                        continue
                    item["type"] = "text"
                    item["content"] = shape.text_frame.text.strip()
                else:
                    continue
                cells.append(item)
        return cells

    @staticmethod
    def get_image_list(cells):
        image_dict = {}
        for cell in cells:
            if cell["type"] == "image":
                uid = uuid.uuid1()
                image_dict[cell["meta"]["id"]] = Image(id=uid, image=cell["content"], image_name=cell["meta"]["name"])
                cell["meta"]["id"] = uid
        return image_dict.values()

    def handle(self, file, pattern_list: List, with_filter: bool, limit: int, get_buffer, save_image):
        try:

            if file.name.endswith(".ppt"):
                cells = self.parse_ppt(file)
            else:
                cells = self.parse_pptx(file)
            image_list = self.get_image_list(cells)
            content = self.to_md(cells)
            if len(image_list) > 0:
                save_image(image_list)
            if pattern_list is not None and len(pattern_list) > 0:
                split_model = SplitModel(pattern_list, with_filter, limit)
            else:
                split_model = SplitModel(default_pattern_list, with_filter=with_filter, limit=limit)
        except BaseException as e:
            traceback.print_exception(e)
            return {'name': file.name,
                    'content': []}
        return {'name': file.name,
                'content': split_model.parse(content)
                }

    def support(self, file, get_buffer):
        file_name: str = file.name.lower()
        if file_name.endswith(".pptx") or file_name.endswith(".ppt"):
            return True
        return False
